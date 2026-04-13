from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "deliverables" / "presentation" / "hra_dashboard_client_presentation.pptx"

TITLE_COLOR = RGBColor(15, 118, 110)
TEXT_COLOR = RGBColor(31, 41, 55)
ACCENT_COLOR = RGBColor(249, 115, 22)
BG_COLOR = RGBColor(245, 243, 239)


SLIDES = [
    {
        "title": "Human Reference Atlas Dashboard",
        "subtitle": "Analysis & Recommendations\nPrepared from the HRA event and request logs",
        "bullets": [],
    },
    {
        "title": "Scope and Dataset",
        "bullets": [
            "Source parquet: 2026-01-13_hra-logs.parquet",
            "Behavioral analysis window: August 5, 2025 to January 12, 2026",
            "111,405 human events, 6,663 sessions, 15 app buckets",
            "Behavioral analysis excludes Bot and AI-Assistant / Bot traffic",
        ],
    },
    {
        "title": "Executive Summary",
        "bullets": [
            "HRA shows real engagement across RUI, EUI, KG Explorer, CDE, and the Portal",
            "Usage is highly concentrated in a minority of sessions",
            "Spatial search is discoverable but has strong funnel drop-off",
            "RUI opacity is a legitimate expert feature, not a mainstream behavior",
            "CDE download behavior is not directly measurable because downloads are not instrumented",
            "Telemetry quality and uneven caching are the biggest cross-cutting issues",
        ],
    },
    {
        "title": "RQ1: Event Frequency",
        "bullets": [
            "Clicks are the dominant event type: 38,961 events",
            "Errors are unusually high: 25,811 events",
            "Median session size is only 3 events",
            "99th percentile session size is about 195 events",
            "A small number of heavy sessions drive a disproportionate share of activity",
        ],
    },
    {
        "title": "RQ2: UI Element Usage",
        "bullets": [
            "Strongest usage appears in RUI registration, EUI results, KG browsing, and CDE setup",
            "2,486 clicked UI paths have fewer than 5 clicks",
            "Several surfaces show heavy hover attention but weak click conversion",
            "This suggests a mix of low-value features and discoverability-to-usability gaps",
        ],
    },
    {
        "title": "RQ3: RUI Opacity",
        "bullets": [
            "206 opacity-related events across 9,430 total RUI events",
            "Opacity is about 2.2% of all RUI activity",
            "Most usage is at the anatomical-structure level, not the global settings panel",
            "Conclusion: useful for advanced users, but not a top-volume feature",
        ],
    },
    {
        "title": "RQ4: EUI Spatial Search",
        "bullets": [
            "1,016 spatial-search events across 69 sessions",
            "Funnel: 66 opened, 50 configured, 27 continued, 24 reviewed results, 7 applied/reset",
            "The feature appears valuable but difficult to complete",
            "Recommendation: improve the configure-to-continue step first",
        ],
    },
    {
        "title": "RQ5: CDE Downloads",
        "bullets": [
            "Users do reach create and visualize stages in CDE",
            "However, no direct histogram or violin download events appear in the logs",
            "This is an instrumentation gap, not proof of no usage",
            "Recommendation: add explicit download events before making product decisions about export behavior",
        ],
    },
    {
        "title": "Performance Findings",
        "bullets": [
            "CDN performs best with 79.03% cache-served requests",
            "KG Explorer page routes perform strongly with 96.83% cache-served requests",
            "CDE page routes are materially better than EUI and RUI page routes",
            "Apps and Portal surfaces show room for meaningful cache and error improvement",
        ],
    },
    {
        "title": "Data-Quality Risks",
        "bullets": [
            "28,087 events have no app attribution and are grouped as Unspecified",
            "25,856 events have no usable UI path",
            "Missing attribution reduces confidence in feature-level rankings and funnel counts",
            "Instrumentation cleanup should be treated as a product and analytics priority",
        ],
    },
    {
        "title": "Recommendations",
        "bullets": [
            "Fix instrumentation first: app, path, session, and CDE download events",
            "Improve EUI spatial-search completion",
            "Review hover-heavy, low-click surfaces for conversion problems",
            "Preserve RUI opacity as an expert feature without over-prioritizing it",
            "Investigate app-page caching for EUI, RUI, Dashboard, and ASCT+B Reporter",
        ],
    },
    {
        "title": "Next Steps",
        "bullets": [
            "Ship instrumentation fixes and rerun the dashboard",
            "Prioritize UX improvements in EUI spatial search",
            "Add direct tracking for CDE chart downloads",
            "Use the next dashboard refresh to validate which product changes actually improve usage and completion",
        ],
    },
]


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR


def style_title(text_frame) -> None:
    paragraph = text_frame.paragraphs[0]
    paragraph.font.name = "Aptos Display"
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = TITLE_COLOR


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.5), Inches(11.2), Inches(0.8))
    style_title(title_box.text_frame)
    title_box.text_frame.text = title

    accent = slide.shapes.add_shape(1, Inches(0.7), Inches(1.25), Inches(1.1), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_COLOR
    accent.line.fill.background()

    body_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.0), Inches(5.1))
    tf = body_box.text_frame
    tf.clear()

    for index, bullet in enumerate(bullets):
        paragraph = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.space_after = Pt(10)
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(22)
        paragraph.font.color.rgb = TEXT_COLOR


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)

    title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.3), Inches(10.8), Inches(1.5))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.font.name = "Aptos Display"
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = TITLE_COLOR

    sub_box = slide.shapes.add_textbox(Inches(0.95), Inches(3.0), Inches(10.4), Inches(1.8))
    sub_tf = sub_box.text_frame
    sub_tf.text = subtitle
    p2 = sub_tf.paragraphs[0]
    p2.font.name = "Aptos"
    p2.font.size = Pt(22)
    p2.font.color.rgb = TEXT_COLOR

    accent = slide.shapes.add_shape(1, Inches(0.95), Inches(2.6), Inches(1.35), Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_COLOR
    accent.line.fill.background()


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    prs.core_properties.title = "Human Reference Atlas Dashboard - Analysis & Recommendations"
    prs.core_properties.subject = "Client presentation based on the HRA dashboard analysis"
    prs.core_properties.author = "OpenAI Codex"

    first = SLIDES[0]
    add_title_slide(prs, first["title"], first["subtitle"])
    for slide in SLIDES[1:]:
        add_bullet_slide(prs, slide["title"], slide["bullets"])

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    print(f"Wrote presentation deck to {OUTPUT}")


if __name__ == "__main__":
    main()
